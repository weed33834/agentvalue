"""深度文档解析服务

对标 RagFlow DeepDoc:
- 表格提取 + 结构化信息 + 版面分析
- 支持文件类型: PDF / DOCX / XLSX / PPTX / TXT / MD
- 解析策略: auto / ocr / structure / hybrid

依赖库 (可能未安装, 用 try/import 降级处理):
- pdfplumber (PDF 文本+表格提取) — 首选
- PyMuPDF / fitz (PDF 图片提取) — 备选
- python-docx (DOCX 解析)
- openpyxl (XLSX 解析)
- python-pptx (PPTX 解析)

异步任务:
- process_task 用 asyncio.create_task() 后台执行
- 解析逻辑本身是同步的 (pdfplumber 等库为同步), 用 asyncio.to_thread 包装

事务边界由路由层控制 (create_task / process_task 内部不 commit)。
"""

from __future__ import annotations

import asyncio
import json
import logging
import os
from datetime import datetime, timezone
from typing import Any, Dict, List, Optional

from sqlalchemy import func, select
from sqlalchemy.ext.asyncio import AsyncSession

from core.database import AsyncSessionLocal
from core.tenant_context import tenant_scope
from models.doc_parsing_models import DocParsingResult, DocParsingTask

logger = logging.getLogger(__name__)

# 支持的文件类型
SUPPORTED_FILE_TYPES = {"pdf", "docx", "xlsx", "pptx", "txt", "md"}

# 支持的解析策略
SUPPORTED_STRATEGIES = {"auto", "ocr", "structure", "hybrid"}

# 任务状态
TASK_STATUS_PENDING = "pending"
TASK_STATUS_PROCESSING = "processing"
TASK_STATUS_COMPLETED = "completed"
TASK_STATUS_FAILED = "failed"


class DocParsingService:
    """深度文档解析服务"""

    def __init__(self, session: AsyncSession):
        self.session = session

    # ===================== 任务 CRUD =====================

    async def create_task(
        self,
        file_path: str,
        file_type: str,
        parse_strategy: str = "auto",
        *,
        tenant_id: str = "default",
    ) -> DocParsingTask:
        """创建解析任务

        Args:
            file_path: 待解析文件路径。
            file_type: 文件类型 (pdf/docx/xlsx/pptx/txt/md)。
            parse_strategy: 解析策略 (auto/ocr/structure/hybrid)。
            tenant_id: 租户 ID。

        Returns:
            创建的 DocParsingTask 对象。

        Raises:
            ValueError: 参数无效或文件不存在。
        """
        if file_type not in SUPPORTED_FILE_TYPES:
            raise ValueError(
                f"不支持的文件类型: {file_type}, 可选: {SUPPORTED_FILE_TYPES}"
            )
        if parse_strategy not in SUPPORTED_STRATEGIES:
            raise ValueError(
                f"不支持的解析策略: {parse_strategy}, 可选: {SUPPORTED_STRATEGIES}"
            )
        if not file_path or not file_path.strip():
            raise ValueError("file_path 不能为空")

        # 路径安全检查: 防止路径遍历 (C2)
        self._validate_file_path(file_path.strip())

        task = DocParsingTask(
            tenant_id=tenant_id,
            file_path=file_path.strip(),
            file_type=file_type,
            parse_strategy=parse_strategy,
            status=TASK_STATUS_PENDING,
        )
        self.session.add(task)
        await self.session.flush()
        logger.info(
            "创建文档解析任务 id=%s file=%s type=%s tenant=%s",
            task.id,
            file_path,
            file_type,
            tenant_id,
        )
        return task

    async def get_task(
        self, task_id: int, *, tenant_id: str = "default"
    ) -> Optional[DocParsingTask]:
        """获取解析任务"""
        return (
            await self.session.execute(
                select(DocParsingTask).where(
                    DocParsingTask.id == task_id,
                    DocParsingTask.tenant_id == tenant_id,
                )
            )
        ).scalar_one_or_none()

    async def list_tasks(
        self,
        *,
        status_filter: Optional[str] = None,
        file_type: Optional[str] = None,
        page: int = 1,
        size: int = 20,
        tenant_id: str = "default",
    ) -> Dict[str, Any]:
        """任务列表 (分页)"""
        base = (
            select(DocParsingTask)
            .where(DocParsingTask.tenant_id == tenant_id)
            .order_by(DocParsingTask.created_at.desc())
        )
        if status_filter:
            base = base.where(DocParsingTask.status == status_filter)
        if file_type:
            base = base.where(DocParsingTask.file_type == file_type)

        total = (
            await self.session.execute(
                select(func.count()).select_from(base.subquery())
            )
        ).scalar() or 0

        offset = (page - 1) * size
        rows = (
            (await self.session.execute(base.offset(offset).limit(size)))
            .scalars()
            .all()
        )

        return {
            "items": [self._task_to_dict(t) for t in rows],
            "total": total,
            "page": page,
            "size": size,
        }

    async def delete_task(self, task_id: int, *, tenant_id: str = "default") -> bool:
        """删除解析任务 (级联删除结果)"""
        task = await self.get_task(task_id, tenant_id=tenant_id)
        if task is None:
            return False
        await self.session.delete(task)
        await self.session.flush()
        return True

    # ===================== 任务执行 =====================

    def schedule_processing(self, task_id: int, *, tenant_id: str = "default") -> None:
        """用 asyncio.create_task() 后台执行解析

        在独立 session 中执行 (避免与请求 session 冲突)。

        Args:
            task_id: 任务 ID。
            tenant_id: 租户 ID (用于设置后台任务租户上下文)。
        """
        asyncio.create_task(self._run_processing(task_id, tenant_id=tenant_id))
        logger.info(
            "调度文档解析任务 task_id=%s tenant=%s (后台执行)", task_id, tenant_id
        )

    async def _run_processing(
        self, task_id: int, *, tenant_id: str = "default"
    ) -> None:
        """后台执行解析 (独立 session)

        使用 tenant_scope 设置租户上下文, 确保后台任务的租户隔离 (M1)。
        """
        with tenant_scope(tenant_id):
            async with AsyncSessionLocal() as session:
                service = DocParsingService(session)
                try:
                    await service.process_task(task_id, tenant_id=tenant_id)
                    await session.commit()
                except Exception as e:
                    logger.exception("文档解析后台任务失败 task_id=%s: %s", task_id, e)
                    await session.rollback()

    async def process_task(
        self, task_id: int, *, tenant_id: str = "default"
    ) -> DocParsingTask:
        """执行解析任务

        1. 加载任务 (按 task_id + tenant_id 过滤, 防跨租户访问), 检查文件是否存在
        2. 路径安全检查 (纵深防御, 防路径遍历)
        3. 标记为 processing
        4. 根据文件类型分发到对应解析器 (同步库用 asyncio.to_thread 包装)
        5. 写入解析结果 (DocParsingResult)
        6. 更新任务统计 (page_count / table_count / image_count) + status=completed

        Args:
            task_id: 任务 ID。
            tenant_id: 租户 ID (用于查询过滤, 防跨租户访问)。

        Returns:
            更新后的 DocParsingTask 对象。

        Raises:
            ValueError: 任务不存在 / 文件不存在 / 解析失败 / 路径不安全。
        """
        task = (
            await self.session.execute(
                select(DocParsingTask).where(
                    DocParsingTask.id == task_id,
                    DocParsingTask.tenant_id == tenant_id,
                )
            )
        ).scalar_one_or_none()
        if task is None:
            raise ValueError(f"解析任务 {task_id} 不存在")

        # 路径安全检查 (纵深防御, C2): 即使 create_task 已检查, 此处再次校验
        self._validate_file_path(task.file_path)

        if not os.path.exists(task.file_path):
            task.status = TASK_STATUS_FAILED
            task.error_message = f"文件不存在: {task.file_path}"
            await self.session.flush()
            raise ValueError(f"文件不存在: {task.file_path}")

        task.status = TASK_STATUS_PROCESSING
        await self.session.flush()

        try:
            # 解析逻辑是同步的, 用 to_thread 包装
            parse_result = await asyncio.to_thread(
                self._dispatch_parse,
                task.file_path,
                task.file_type,
                task.parse_strategy,
            )

            # 写入解析结果
            page_count = 0
            table_count = 0
            image_count = 0
            for item in parse_result.get("results", []):
                result = DocParsingResult(
                    tenant_id=task.tenant_id,
                    task_id=task.id,
                    page_num=item.get("page_num"),
                    content_type=item.get("content_type", "text"),
                    content=item.get("content", ""),
                    bounding_box=item.get("bounding_box"),
                    metadata_=item.get("metadata"),
                )
                self.session.add(result)
                if item.get("content_type") == "table":
                    table_count += 1
                elif item.get("content_type") == "image":
                    image_count += 1

            page_count = parse_result.get("page_count", 0)
            # 若未统计到 table/image, 从 parse_result 元数据补充
            table_count = table_count or parse_result.get("table_count", 0)
            image_count = image_count or parse_result.get("image_count", 0)

            task.page_count = page_count
            task.table_count = table_count
            task.image_count = image_count
            task.result = {
                "page_count": page_count,
                "table_count": table_count,
                "image_count": image_count,
                "content_count": len(parse_result.get("results", [])),
            }
            task.status = TASK_STATUS_COMPLETED
            task.completed_at = datetime.now(timezone.utc)
            task.error_message = None
            await self.session.flush()
            logger.info(
                "文档解析完成 task_id=%s pages=%s tables=%s images=%s",
                task_id,
                page_count,
                table_count,
                image_count,
            )
        except Exception as e:
            logger.exception("文档解析失败 task_id=%s: %s", task_id, e)
            task.status = TASK_STATUS_FAILED
            task.error_message = str(e)
            await self.session.flush()
            raise

        return task

    # ===================== 路径安全 =====================

    @staticmethod
    def _get_allowed_root() -> str:
        """获取允许的文件根目录 (从 settings.attachment_dir)"""
        from core.config import get_settings

        return os.path.realpath(get_settings().attachment_dir)

    def _validate_file_path(self, file_path: str) -> str:
        """验证文件路径在允许的根目录内, 防止路径遍历 (C2)

        用 os.path.realpath 解析路径 (消除 .. 和符号链接),
        然后验证解析后的绝对路径是否在允许的根目录内。

        Args:
            file_path: 待验证的文件路径。

        Returns:
            解析后的绝对路径。

        Raises:
            ValueError: 路径不在允许的目录内。
        """
        allowed_root = self._get_allowed_root()
        real_path = os.path.realpath(file_path)

        # 验证解析后的路径在允许的根目录内
        if real_path != allowed_root and not real_path.startswith(
            allowed_root + os.sep
        ):
            raise ValueError(
                f"文件路径不在允许的目录内: {file_path} " f"(允许目录: {allowed_root})"
            )
        return real_path

    # ===================== 结果查询 =====================

    async def get_task_results(
        self,
        task_id: int,
        page_num: Optional[int] = None,
        *,
        tenant_id: str = "default",
    ) -> List[DocParsingResult]:
        """获取解析结果 (可按页过滤)"""
        base = select(DocParsingResult).where(
            DocParsingResult.task_id == task_id,
            DocParsingResult.tenant_id == tenant_id,
        )
        if page_num is not None:
            base = base.where(DocParsingResult.page_num == page_num)
        base = base.order_by(DocParsingResult.page_num, DocParsingResult.id)
        result = await self.session.execute(base)
        return list(result.scalars().all())

    # ===================== 解析分发 =====================

    def _dispatch_parse(
        self, file_path: str, file_type: str, strategy: str
    ) -> Dict[str, Any]:
        """根据文件类型分发到对应解析器

        Returns:
            {
                "results": [{page_num, content_type, content, bounding_box, metadata}],
                "page_count": int,
                "table_count": int,
                "image_count": int,
            }
        """
        if file_type == "pdf":
            return self._parse_pdf(file_path, strategy)
        elif file_type == "docx":
            return self._parse_docx(file_path)
        elif file_type == "xlsx":
            return self._parse_xlsx(file_path)
        elif file_type in ("txt", "md"):
            return self._parse_generic(file_path)
        elif file_type == "pptx":
            return self._parse_pptx(file_path)
        else:
            raise ValueError(f"不支持的文件类型: {file_type}")

    # ===================== PDF 解析 =====================

    def _parse_pdf(self, file_path: str, strategy: str) -> Dict[str, Any]:
        """PDF 解析 (用 pdfplumber 提取文本/表格, PyMuPDF 提取图片)

        解析内容:
        - 文本: 按页提取文本段落
        - 表格: 转为结构化 JSON (含行列数据)
        - 图片: 记录图片引用 (bbox + page)
        - 版面: 保留 bounding_box 用于版面分析

        Args:
            file_path: PDF 文件路径。
            strategy: 解析策略 (auto/ocr/structure/hybrid, 当前实现 structure 模式)。

        Raises:
            ValueError: pdfplumber 未安装。
        """
        try:
            import pdfplumber  # type: ignore
        except ImportError:
            raise ValueError("pdfplumber 未安装, 请运行 pip install pdfplumber 后重试")

        results: List[Dict[str, Any]] = []
        page_count = 0
        table_count = 0
        image_count = 0

        with pdfplumber.open(file_path) as pdf:
            page_count = len(pdf.pages)
            for page_idx, page in enumerate(pdf.pages, start=1):
                # 1. 提取文本
                page_text = page.extract_text() or ""
                if page_text.strip():
                    # 按段落分割 (空行分隔)
                    paragraphs = [
                        p.strip() for p in page_text.split("\n\n") if p.strip()
                    ]
                    if not paragraphs:
                        paragraphs = [page_text.strip()]
                    for para in paragraphs:
                        results.append(
                            {
                                "page_num": page_idx,
                                "content_type": "text",
                                "content": para,
                                "bounding_box": None,
                                "metadata": {"char_count": len(para)},
                            }
                        )

                # 2. 提取表格
                tables = self._extract_tables(page)
                for table in tables:
                    results.append(
                        {
                            "page_num": page_idx,
                            "content_type": "table",
                            "content": json.dumps(table, ensure_ascii=False),
                            "bounding_box": None,
                            "metadata": {
                                "rows": len(table.get("rows", [])),
                                "columns": len(table.get("headers", [])),
                            },
                        }
                    )
                    table_count += 1

                # 3. 提取图片 (用 PyMuPDF 备选, 或记录 image bbox)
                images = page.images or []
                for img in images:
                    results.append(
                        {
                            "page_num": page_idx,
                            "content_type": "image",
                            "content": f"image_page{page_idx}_{image_count}",
                            "bounding_box": {
                                "x0": float(img.get("x0", 0)),
                                "y0": float(img.get("top", 0)),
                                "x1": float(img.get("x1", 0)),
                                "y1": float(img.get("bottom", 0)),
                            },
                            "metadata": {
                                "width": float(img.get("width", 0)),
                                "height": float(img.get("height", 0)),
                            },
                        }
                    )
                    image_count += 1

        return {
            "results": results,
            "page_count": page_count,
            "table_count": table_count,
            "image_count": image_count,
        }

    def _extract_tables(self, pdf_page: Any) -> List[Dict[str, Any]]:
        """从 PDF 页面提取表格, 转为结构化 JSON

        每个表格转为:
        {
            "headers": [col1, col2, ...],
            "rows": [[v1, v2, ...], ...],
            "row_count": N,
            "col_count": M,
        }

        Args:
            pdf_page: pdfplumber Page 对象。

        Returns:
            表格列表 (结构化 JSON)。
        """
        tables_out: List[Dict[str, Any]] = []
        try:
            tables = pdf_page.extract_tables() or []
            for table in tables:
                if not table or len(table) < 1:
                    continue
                # 第一行作为表头
                headers = [
                    str(c).strip() if c else f"col_{i}" for i, c in enumerate(table[0])
                ]
                rows = []
                for row in table[1:]:
                    rows.append([str(c).strip() if c is not None else "" for c in row])
                tables_out.append(
                    {
                        "headers": headers,
                        "rows": rows,
                        "row_count": len(rows),
                        "col_count": len(headers),
                    }
                )
        except Exception as e:
            logger.warning("表格提取失败: %s", e)
        return tables_out

    # ===================== DOCX 解析 =====================

    def _parse_docx(self, file_path: str) -> Dict[str, Any]:
        """DOCX 解析 (用 python-docx)

        解析内容:
        - 标题 (heading): 按 style 判断层级
        - 段落 (text): 普通段落
        - 表格 (table): 转 JSON
        - 列表 (list): 按样式判断

        Raises:
            ValueError: python-docx 未安装。
        """
        try:
            from docx import Document  # type: ignore
        except ImportError:
            raise ValueError(
                "python-docx 未安装, 请运行 pip install python-docx 后重试"
            )

        results: List[Dict[str, Any]] = []
        doc = Document(file_path)

        # 遍历 body 元素 (段落 + 表格, 保持顺序)
        from docx.table import Table  # type: ignore
        from docx.text.paragraph import Paragraph  # type: ignore

        page_num = 1  # DOCX 无页码概念, 统一记为 1
        table_count = 0

        body = doc.element.body
        para_idx = 0
        table_idx = 0
        paragraphs = doc.paragraphs
        tables = doc.tables

        for child in body.iterchildren():
            if child.tag.endswith("}p"):
                if para_idx < len(paragraphs):
                    para = paragraphs[para_idx]
                    para_idx += 1
                    text = (para.text or "").strip()
                    if not text:
                        continue
                    style_name = (para.style.name or "").lower() if para.style else ""
                    if "heading" in style_name or "title" in style_name:
                        # 提取标题层级
                        level = 1
                        for ch in style_name:
                            if ch.isdigit():
                                level = int(ch)
                                break
                        results.append(
                            {
                                "page_num": page_num,
                                "content_type": "heading",
                                "content": text,
                                "bounding_box": None,
                                "metadata": {"level": level, "style": style_name},
                            }
                        )
                    elif "list" in style_name:
                        results.append(
                            {
                                "page_num": page_num,
                                "content_type": "list",
                                "content": text,
                                "bounding_box": None,
                                "metadata": {"style": style_name},
                            }
                        )
                    else:
                        results.append(
                            {
                                "page_num": page_num,
                                "content_type": "text",
                                "content": text,
                                "bounding_box": None,
                                "metadata": {"char_count": len(text)},
                            }
                        )
            elif child.tag.endswith("}tbl"):
                if table_idx < len(tables):
                    table = tables[table_idx]
                    table_idx += 1
                    rows_data = []
                    for row in table.rows:
                        rows_data.append([cell.text.strip() for cell in row.cells])
                    headers = rows_data[0] if rows_data else []
                    table_json = {
                        "headers": headers,
                        "rows": rows_data[1:] if len(rows_data) > 1 else [],
                        "row_count": max(0, len(rows_data) - 1),
                        "col_count": len(headers),
                    }
                    results.append(
                        {
                            "page_num": page_num,
                            "content_type": "table",
                            "content": json.dumps(table_json, ensure_ascii=False),
                            "bounding_box": None,
                            "metadata": {
                                "rows": table_json["row_count"],
                                "columns": table_json["col_count"],
                            },
                        }
                    )
                    table_count += 1

        return {
            "results": results,
            "page_count": 1,
            "table_count": table_count,
            "image_count": 0,
        }

    # ===================== XLSX 解析 =====================

    def _parse_xlsx(self, file_path: str) -> Dict[str, Any]:
        """Excel 解析 (用 openpyxl)

        每个 sheet 作为一个表格结果 (含表头 + 行数据)。

        Raises:
            ValueError: openpyxl 未安装。
        """
        try:
            from openpyxl import load_workbook  # type: ignore
        except ImportError:
            raise ValueError("openpyxl 未安装, 请运行 pip install openpyxl 后重试")

        results: List[Dict[str, Any]] = []
        table_count = 0

        wb = load_workbook(file_path, data_only=True, read_only=True)
        for sheet_idx, sheet in enumerate(wb.worksheets, start=1):
            rows_iter = sheet.iter_rows(values_only=True)
            rows_list = list(rows_iter)
            if not rows_list:
                continue
            headers = [
                str(c).strip() if c is not None else f"col_{i}"
                for i, c in enumerate(rows_list[0])
            ]
            data_rows = [
                [str(c).strip() if c is not None else "" for c in row]
                for row in rows_list[1:]
            ]
            table_json = {
                "headers": headers,
                "rows": data_rows,
                "row_count": len(data_rows),
                "col_count": len(headers),
                "sheet_name": sheet.title,
            }
            results.append(
                {
                    "page_num": sheet_idx,
                    "content_type": "table",
                    "content": json.dumps(table_json, ensure_ascii=False),
                    "bounding_box": None,
                    "metadata": {
                        "sheet_name": sheet.title,
                        "rows": len(data_rows),
                        "columns": len(headers),
                    },
                }
            )
            table_count += 1
        wb.close()

        return {
            "results": results,
            "page_count": len(wb.worksheets),
            "table_count": table_count,
            "image_count": 0,
        }

    # ===================== PPTX 解析 =====================

    def _parse_pptx(self, file_path: str) -> Dict[str, Any]:
        """PPTX 解析 (用 python-pptx)

        每张幻灯片作为一页, 提取标题 + 文本框 + 表格。

        Raises:
            ValueError: python-pptx 未安装。
        """
        try:
            from pptx import Presentation  # type: ignore
        except ImportError:
            raise ValueError(
                "python-pptx 未安装, 请运行 pip install python-pptx 后重试"
            )

        results: List[Dict[str, Any]] = []
        table_count = 0
        image_count = 0

        prs = Presentation(file_path)
        for slide_idx, slide in enumerate(prs.slides, start=1):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        # 第一张幻灯片的标题作为 heading
                        is_title = shape == slide.shapes.title
                        results.append(
                            {
                                "page_num": slide_idx,
                                "content_type": "heading" if is_title else "text",
                                "content": text,
                                "bounding_box": None,
                                "metadata": {"shape_name": shape.name},
                            }
                        )
                if shape.has_table:
                    tbl = shape.table
                    rows_data = []
                    for row in tbl.rows:
                        rows_data.append([cell.text.strip() for cell in row.cells])
                    headers = rows_data[0] if rows_data else []
                    table_json = {
                        "headers": headers,
                        "rows": rows_data[1:] if len(rows_data) > 1 else [],
                        "row_count": max(0, len(rows_data) - 1),
                        "col_count": len(headers),
                    }
                    results.append(
                        {
                            "page_num": slide_idx,
                            "content_type": "table",
                            "content": json.dumps(table_json, ensure_ascii=False),
                            "bounding_box": None,
                            "metadata": {
                                "rows": table_json["row_count"],
                                "columns": table_json["col_count"],
                            },
                        }
                    )
                    table_count += 1
                if shape.shape_type == 13:  # PICTURE
                    results.append(
                        {
                            "page_num": slide_idx,
                            "content_type": "image",
                            "content": f"image_slide{slide_idx}_{image_count}",
                            "bounding_box": None,
                            "metadata": {
                                "shape_name": shape.name,
                                "width": int(shape.width) if shape.width else 0,
                                "height": int(shape.height) if shape.height else 0,
                            },
                        }
                    )
                    image_count += 1

        return {
            "results": results,
            "page_count": len(prs.slides),
            "table_count": table_count,
            "image_count": image_count,
        }

    # ===================== 通用文本解析 =====================

    def _parse_generic(self, file_path: str) -> Dict[str, Any]:
        """通用文本解析 (TXT / MD)

        按段落分割 (空行分隔), Markdown 标题 (#) 识别为 heading。
        """
        results: List[Dict[str, Any]] = []
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                content = f.read()
        except UnicodeDecodeError:
            with open(file_path, "r", encoding="gbk", errors="ignore") as f:
                content = f.read()

        is_markdown = file_path.lower().endswith(".md")
        paragraphs = [p.strip() for p in content.split("\n\n") if p.strip()]
        if not paragraphs:
            paragraphs = [line.strip() for line in content.split("\n") if line.strip()]

        for para in paragraphs:
            if is_markdown and para.startswith("#"):
                level = len(para) - len(para.lstrip("#"))
                results.append(
                    {
                        "page_num": 1,
                        "content_type": "heading",
                        "content": para.lstrip("# ").strip(),
                        "bounding_box": None,
                        "metadata": {"level": level, "format": "markdown"},
                    }
                )
            else:
                results.append(
                    {
                        "page_num": 1,
                        "content_type": "text",
                        "content": para,
                        "bounding_box": None,
                        "metadata": {"char_count": len(para)},
                    }
                )

        return {
            "results": results,
            "page_count": 1,
            "table_count": 0,
            "image_count": 0,
        }

    # ===================== 元数据提取 =====================

    def _extract_metadata(self, file_path: str, file_type: str) -> Dict[str, Any]:
        """提取文件元数据

        Returns:
            {
                "file_name": str,
                "file_size": int (bytes),
                "file_type": str,
                "modified_time": float (timestamp),
            }
        """
        stat = os.stat(file_path)
        return {
            "file_name": os.path.basename(file_path),
            "file_size": stat.st_size,
            "file_type": file_type,
            "modified_time": stat.st_mtime,
        }

    # ===================== 序列化 =====================

    @staticmethod
    def _task_to_dict(t: DocParsingTask) -> Dict[str, Any]:
        """DocParsingTask → dict"""
        return {
            "id": t.id,
            "tenant_id": t.tenant_id,
            "file_path": t.file_path,
            "file_type": t.file_type,
            "parse_strategy": t.parse_strategy,
            "status": t.status,
            "result": t.result if isinstance(t.result, dict) else None,
            "error_message": t.error_message,
            "page_count": t.page_count,
            "table_count": t.table_count,
            "image_count": t.image_count,
            "created_at": t.created_at.isoformat() if t.created_at else None,
            "completed_at": t.completed_at.isoformat() if t.completed_at else None,
        }

    @staticmethod
    def _result_to_dict(r: DocParsingResult) -> Dict[str, Any]:
        """DocParsingResult → dict"""
        return {
            "id": r.id,
            "tenant_id": r.tenant_id,
            "task_id": r.task_id,
            "page_num": r.page_num,
            "content_type": r.content_type,
            "content": r.content,
            "bounding_box": (
                r.bounding_box if isinstance(r.bounding_box, dict) else None
            ),
            "metadata": r.metadata_ if isinstance(r.metadata_, dict) else None,
            "created_at": r.created_at.isoformat() if r.created_at else None,
        }
